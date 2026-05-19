from __future__ import annotations

import hashlib
import json
from collections.abc import Callable
from dataclasses import dataclass
from pathlib import Path

from app.core.config import get_settings
from app.services.text_utils import SectionDraft, build_sections_from_plain_text, normalize_text, split_paragraphs

SUPPORTED_EXTENSIONS = {
    ".md",
    ".txt",
    ".docx",
    ".pdf",
    ".html",
    ".htm",
    ".xlsx",
    ".pptx",
    ".png",
    ".jpg",
    ".jpeg",
    ".webp",
    ".tif",
    ".tiff",
}

IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".webp", ".tif", ".tiff"}

ProgressNote = Callable[[str], None]


@dataclass
class ParseOutcome:
    sections: list[SectionDraft]
    ocr_meta: str


def _empty_outcome(sections: list[SectionDraft]) -> ParseOutcome:
    return ParseOutcome(sections=sections, ocr_meta="")


def _content_hash(path: Path, raw_content_hash: str | None) -> str:
    if raw_content_hash:
        return raw_content_hash
    return hashlib.sha256(path.read_bytes()).hexdigest()


def parse_document(
    path: Path,
    *,
    raw_content_hash: str | None = None,
    progress_note: ProgressNote | None = None,
) -> ParseOutcome:
    suffix = path.suffix.lower()
    if suffix not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"暂不支持的文件格式：{suffix}")
    if suffix == ".md":
        return _empty_outcome(parse_markdown(path))
    if suffix == ".txt":
        return _empty_outcome(build_sections_from_plain_text(path.read_text(encoding="utf-8"), path.stem))
    if suffix == ".docx":
        return _empty_outcome(parse_docx(path))
    if suffix == ".pdf":
        return parse_pdf(path, raw_content_hash=raw_content_hash, progress_note=progress_note)
    if suffix in (".html", ".htm"):
        return _empty_outcome(parse_html_file(path))
    if suffix == ".xlsx":
        return _empty_outcome(parse_xlsx(path))
    if suffix == ".pptx":
        return _empty_outcome(parse_pptx(path))
    if suffix in IMAGE_SUFFIXES:
        return parse_image(path, raw_content_hash=_content_hash(path, raw_content_hash), progress_note=progress_note)
    raise ValueError(f"暂不支持的文件格式：{suffix}")


def parse_markdown(path: Path) -> list[SectionDraft]:
    lines = path.read_text(encoding="utf-8").splitlines()
    sections: list[SectionDraft] = []
    current = SectionDraft(title=path.stem, level=1, paragraphs=[])
    buffer: list[str] = []

    def flush_paragraphs() -> None:
        nonlocal buffer
        if buffer:
            current.paragraphs.extend(split_paragraphs("\n".join(buffer)))
            buffer = []

    for line in lines:
        if line.startswith("#"):
            marker, _, title = line.partition(" ")
            if marker and set(marker) == {"#"} and title.strip():
                flush_paragraphs()
                if current.paragraphs or not sections:
                    sections.append(current)
                current = SectionDraft(title=title.strip(), level=len(marker), paragraphs=[])
                continue
        buffer.append(line)

    flush_paragraphs()
    if current.paragraphs or not sections:
        sections.append(current)
    return sections


def parse_docx(path: Path) -> list[SectionDraft]:
    try:
        from docx import Document
    except ImportError as exc:
        raise RuntimeError("解析 DOCX 需要安装 python-docx") from exc

    doc = Document(path)
    sections: list[SectionDraft] = []
    current = SectionDraft(title=path.stem, level=1, paragraphs=[])

    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        style_name = paragraph.style.name.lower() if paragraph.style else ""
        if style_name.startswith("heading") or style_name.startswith("标题"):
            if current.paragraphs or not sections:
                sections.append(current)
            level = 2
            for token in style_name.split():
                if token.isdigit():
                    level = int(token)
                    break
            current = SectionDraft(title=text, level=level, paragraphs=[])
        else:
            current.paragraphs.append(text)

    if current.paragraphs or not sections:
        sections.append(current)
    return sections


def parse_html_file(path: Path) -> list[SectionDraft]:
    raw = path.read_text(encoding="utf-8", errors="replace")
    return parse_html_content(raw, path.stem)


def parse_html_content(html: str, fallback_title: str) -> list[SectionDraft]:
    try:
        from bs4 import BeautifulSoup
    except ImportError as exc:
        raise RuntimeError("解析 HTML 需要安装 beautifulsoup4、lxml") from exc

    soup = BeautifulSoup(html, "lxml")
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    title_el = soup.find("title")
    page_title = title_el.get_text(strip=True) if title_el else ""
    text = soup.get_text(separator="\n")
    text = normalize_text(text)
    if not text.strip():
        raise ValueError("HTML 未提取到正文文本")
    heading = page_title or fallback_title
    return build_sections_from_plain_text(text, heading)


def parse_xlsx(path: Path) -> list[SectionDraft]:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise RuntimeError("解析 XLSX 需要安装 openpyxl") from exc

    workbook = load_workbook(path, read_only=True, data_only=True)
    sections: list[SectionDraft] = []
    try:
        for sheet in workbook.worksheets:
            lines: list[str] = []
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
                if cells:
                    lines.append("\t".join(cells))
            block = normalize_text("\n".join(lines))
            paragraphs = split_paragraphs(block) if block else []
            if not paragraphs and block:
                paragraphs = [block]
            title = f"工作表: {sheet.title}"
            if paragraphs:
                sections.append(SectionDraft(title=title, level=2, paragraphs=paragraphs))
    finally:
        workbook.close()
    if not sections:
        raise ValueError("表格中未读取到文本内容")
    return sections


def parse_pptx(path: Path) -> list[SectionDraft]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("解析 PPTX 需要安装 python-pptx") from exc

    prs = Presentation(str(path))
    sections: list[SectionDraft] = []
    for index, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                t = (shape.text or "").strip()
                if t:
                    texts.append(t)
        body = normalize_text("\n".join(texts))
        if not body:
            continue
        paras = split_paragraphs(body)
        title = f"幻灯片 {index}"
        sections.append(SectionDraft(title=title, level=2, paragraphs=paras))
    if not sections:
        raise ValueError("演示文稿中未提取到文本")
    return sections


def parse_image(
    path: Path,
    *,
    raw_content_hash: str,
    progress_note: ProgressNote | None = None,
) -> ParseOutcome:
    from app.services.ocr_engine import ocr_image_path

    text, meta = ocr_image_path(path, raw_content_hash=raw_content_hash, progress=progress_note)
    text = normalize_text(text)
    if not text.strip():
        raise ValueError("OCR 未从图片中识别到文本，请检查语言包或图像清晰度")
    sections = build_sections_from_plain_text(text, path.stem)
    return ParseOutcome(sections=sections, ocr_meta=json.dumps(meta, ensure_ascii=False))


def parse_pdf(
    path: Path,
    *,
    raw_content_hash: str | None = None,
    progress_note: ProgressNote | None = None,
) -> ParseOutcome:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise RuntimeError("解析 PDF 需要安装 pypdf") from exc

    reader = PdfReader(str(path))
    page_texts: list[str] = []
    for page in reader.pages:
        page_texts.append(page.extract_text() or "")
    text = normalize_text("\n\n".join(page_texts))
    settings = get_settings()
    threshold = max(0, settings.ocr_pdf_min_text_chars)

    def _cjk_ratio(s: str) -> float:
        compact = "".join(s.split())
        if not compact:
            return 0.0
        n = sum(1 for ch in compact if "\u4e00" <= ch <= "\u9fff" or "\u3400" <= ch <= "\u4dbf")
        return n / len(compact)

    long_enough = len(text.strip()) >= threshold and bool(text.strip())
    min_cjk = float(settings.ocr_pdf_min_cjk_ratio or 0.0)
    suspicious_layer = long_enough and min_cjk > 0.0 and _cjk_ratio(text) < min_cjk

    use_text_layer = (
        not settings.ocr_pdf_force_visual
        and long_enough
        and not suspicious_layer
    )

    if use_text_layer:
        return _empty_outcome(build_sections_from_plain_text(text, path.stem))

    if suspicious_layer and not settings.ocr_enabled:
        raise ValueError(
            "PDF 文字层字符数足够，但汉字占比过低，疑似编码/字体映射异常（常见于部分国标电子版）。"
            "请在 .env 开启 OCR_ENABLED=true，或设置 OCR_PDF_FORCE_VISUAL=true 强制按页渲染识别。"
        )

    if not settings.ocr_enabled:
        if settings.ocr_pdf_force_visual:
            raise ValueError(
                "已设置 OCR_PDF_FORCE_VISUAL=true，须同时 OCR_ENABLED=true，并安装 uv sync --extra ocr 与本机 Tesseract。"
            )
        raise ValueError(
            "PDF 未提取到足够文本（可能为扫描件）。请在 .env 设置 OCR_ENABLED=true，"
            "执行 uv sync --extra ocr，并安装本机 Tesseract 与语言包。"
        )

    from app.services.ocr_engine import ocr_pdf_scanned

    h = _content_hash(path, raw_content_hash)
    text2, meta = ocr_pdf_scanned(path, raw_content_hash=h, progress=progress_note)
    text2 = normalize_text(text2)
    if not text2.strip():
        raise ValueError("OCR 未从 PDF 中识别到文本，请检查语言包或提高 OCR_PDF_DPI")
    sections = build_sections_from_plain_text(text2, path.stem)
    payload = {
        **meta,
        "source": "pdf_scan",
        "text_layer_chars": len(text.strip()),
        "threshold": threshold,
        "used_text_layer": False,
        "force_visual": settings.ocr_pdf_force_visual,
        "suspicious_text_layer": suspicious_layer,
    }
    return ParseOutcome(sections=sections, ocr_meta=json.dumps(payload, ensure_ascii=False))
